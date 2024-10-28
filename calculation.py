# Imports de las librerías necesarias
import time
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from matplotlib.ticker import FuncFormatter
from wordcloud import WordCloud, STOPWORDS, ImageColorGenerator
import textwrap
import os
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Definición de las constantes y configuraciones
social_networks_prensa = ['@EdeesteRD', '@edeesterd', "edenorterd", "edesurrd", '@precisionportal', 'panoramatv_', '@mananerord', '@noticiaspimentl', '@MIC_RD', "siegobrd", '@soldelamananard', '@drmariolama', '@labazucaenlared', '@eldia_do', '@diariosaluddo', '@noticiasrnn', '@periodicohoy', '@telemicrohd', '@precisionportal', '@telenoticiasrd', '@sin24horas', 'sin24horas', '@diariolibre', 'diariolibre', '@cdn37', 'cdn.com.do', '@listindiario', 'listindiario', 'bavaronews', '@acentodiario', '@anoticias7', '@z101digital', 'z101digital', 'revista110', 'rcavada', '@rcavada', 'robertocavada', 'ndigital', 'bavarodigital', '@elnuevodiariord', '@acentodiario', '@deultimominutomedia', '@aeromundo_ggomez', '@cachichatv', '@remolachanews', '@ntelemicro5', '@megadiariord', '@noticiasmn', '@eldemocratard', '@elpregonerord', '@infoturdom', '@comunidadojala', '@panorama_do', '@panoramard1', '@invertix', '@derechoasaberlo', '@rccmediard', '@n16noticias', '@lospueblistas']

social_network_sources = {
    'Twitter': 'Redes Sociales',
    'Youtube': 'Redes Sociales',
    'Instagram': 'Redes Sociales',
    'Facebook': 'Redes Sociales',
    'Pinterest': 'Redes Sociales',
    'Reddit': 'Redes Sociales',
    'TikTok': 'Redes Sociales',
    'Twitch': 'Redes Sociales',
}

# Función para formatear el alcance
def format_number(number):
    """Format a number to display in millions (M) or thousands (k)"""
    if number >= 1_000_000:
        return f"{number / 1_000_000:.1f}M"
    elif number >= 1_000:
        return f"{number / 1_000:.1f}k"
    return str(number)

# Funcion para formatear el alcance solo al crear las tablas en las funciones top_influencers_prensa_digital(), top_influencers_redes_sociales_by_posts() y top_influencers_redes_sociales_by_reach()
def format_reach(number):
    if pd.isna(number):
        return number
    else:
        return f"{int(number):,}"
    
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def set_text_style(shape, text, font_name='Effra', font_size=Pt(14), center=True):
    if shape.has_text_frame:
        # Obtiene el primer párrafo
        p = shape.text_frame.paragraphs[0] if shape.text_frame.paragraphs else shape.text_frame.add_paragraph()
        
        # Limpia el texto existente
        p.clear()  # Elimina el texto actual en el párrafo
        
        # Establece el nuevo texto y formato
        p.text = text
        p.font.name = font_name
        p.font.size = font_size
        p.font.color.rgb = RGBColor(0, 0, 0)  # Color negro

        if center:
            p.alignment = PP_ALIGN.CENTER  # Centra el texto
        else:
            p.alignment = PP_ALIGN.LEFT  # Alinea a la izquierda


# Función para cargar y limpiar el DataFrame
def load_and_clean_data(file_path):
    df = pd.read_csv(file_path, encoding='utf-16', sep='\t')
    columns_to_delete = [
        'Opening Text', 'Subregion', 'Desktop Reach', 'Mobile Reach',
        'Twitter Social Echo', 'Facebook Social Echo', 'Reddit Social Echo',
        'National Viewership', 'AVE', 'State', 'City',
        'Social Echo Total', 'Editorial Echo', 'Views', 'Estimated Views',
        'Likes', 'Replies', 'Retweets', 'Comments', 'Shares', 'Reactions',
        'Threads', 'Is Verified'
    ]
    df_cleaned = df.drop(columns=columns_to_delete, errors='ignore')
    df_cleaned['Hit Sentence'] = df_cleaned['Headline'].where(~df_cleaned['Headline'].isna(), df_cleaned['Hit Sentence'])
    exclude_sources = ['Blogs', 'Twitter', 'Youtube', 'Instagram', 'Facebook', 'Pinterest', 'Reddit', 'TikTok', 'Twitch']
    mask = ~df_cleaned['Source'].isin(exclude_sources)
    df_cleaned.loc[mask, 'Influencer'] = df_cleaned.loc[mask, 'Source']
    df_cleaned['Plataforma'] = df_cleaned['Source'].apply(lambda x: social_network_sources.get(x, 'Prensa Digital'))
    return df_cleaned

# Función para actualizar la columna 'Influencer' para comentarios en Facebook
def update_influencer(row):
    if row['Source'] == 'Facebook' and row['Reach'] == 0:
        return "Comment on " + row['Influencer']
    return row['Influencer']

# Función para actualizar el sentimiento
def update_sentiment(row):
    sentiment = row['Sentiment']
    # Cambia "Unknown" o NaN a "Neutral"
    if sentiment == "Unknown" or pd.isna(sentiment):
        return "Neutral"
    return sentiment

# Función para calcular métricas de resumen
def calculate_summary_metrics(df):
    total_mentions = len(df)
    count_of_authors = df['Influencer'].nunique()
    estimated_reach = df.groupby('Influencer')['Reach'].max().sum()
    formatted_estimated_reach = format_number(estimated_reach)
    return total_mentions, count_of_authors, formatted_estimated_reach

# Función para generar gráfico de evolución temporal de menciones
def create_mentions_evolution_chart(df, date_column='Alternate Date Format', time_column='Time', output_path='scratch/convEvolution.png'):
    df['date'] = pd.to_datetime(df[date_column], format='%d-%b-%y').dt.date
    df['hour'] = pd.to_datetime(df[time_column], format='%I:%M %p').dt.strftime('%I %p')
    df_grouped = df.groupby(['date', 'hour']).size().reset_index(name='count')
    df_grouped['datetime'] = pd.to_datetime(df_grouped['date'].astype(str) + ' ' + df_grouped['hour'])
    df_grouped = df_grouped.sort_values('datetime')
    
    fig, ax = plt.subplots(figsize=(13, 7))
    ax.plot(df_grouped['datetime'], df_grouped['count'], linewidth=3.5, color='orange')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.xaxis.set_major_locator(mdates.HourLocator(interval=2))
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%d-%b %I %p'))
    plt.xticks(rotation=90, ha='right')
    plt.tight_layout()
    plt.savefig(output_path, transparent=True)
    
def create_mentions_evolution_chart_by_date(df, date_column='Alternate Date Format', output_path='scratch/convEvolution.png'):
    df['date'] = pd.to_datetime(df[date_column], format='%d-%b-%y').dt.date
    df_grouped = df.groupby('date').size().reset_index(name='count')
    df_grouped = df_grouped.sort_values('date')
    
    fig, ax = plt.subplots(figsize=(13, 7))
    ax.plot(df_grouped['date'], df_grouped['count'], linewidth=3.5, color='orange')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.xaxis.set_major_locator(mdates.DayLocator(interval=1))
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%d-%b'))
    plt.xticks(rotation=90, ha='right')
    plt.tight_layout()
    plt.savefig(output_path, transparent=True)

# Función para generar gráfico de pie de sentimientos
def create_sentiment_pie_chart(df, output_path='scratch/sentiment_pie_chart.png'):
    sentiment_counts = df[df['Sentiment'] != 'Not Rated']['Sentiment'].value_counts()
    sentiment_colors = {'Negative': '#ff0000', 'Positive': '#00b050', 'Neutral': '#BFBFBF'}
    labels = sentiment_counts.index.to_list()
    sizes = sentiment_counts.to_list()
    
    plt.figure(figsize=(10, 10))
    colors = [sentiment_colors.get(label, 'lightgray') for label in labels]
    plt.pie(sizes, autopct="%1.1f%%", colors=colors, textprops=dict(backgroundcolor='w', fontsize=14, fontweight='bold'), startangle=140)
    plt.axis('equal')
    plt.savefig(output_path, transparent=True)

# Función para distribución por plataforma
def distribucion_plataforma(df):
    platform_counts = df['Plataforma'].value_counts()
    max_reach_per_platform = df.groupby('Plataforma')['Reach'].max().apply(format_number)
    return platform_counts.to_dict(), max_reach_per_platform.to_dict()

# Función para extraer las 5 'Hit Sentence' principales
def get_top_hit_sentences(df):
    top_influencers = df[df['Plataforma'] == "Prensa Digital"].sort_values(by=['Reach'], ascending=False).head(5)['Hit Sentence']
    return [textwrap.fill(sentence, width=100) for sentence in top_influencers]

# Funciones para generar las tablas de Top Users
def top_influencers_prensa_digital(df_cleaned):
    filter_prensa_digital = df_cleaned['Plataforma'] == "Prensa Digital"
    df_prensa = df_cleaned[filter_prensa_digital]
    df_prensa_grouped = df_prensa.groupby('Influencer')[['Reach']].agg(['count', 'max'])
    df_prensa_grouped.columns = ['Posts', 'Max Reach']
    df_prensa_grouped = df_prensa_grouped.sort_values(by='Posts', ascending=False).head(10)
    df_prensa_grouped.reset_index(inplace=True)
    df_prensa_grouped['Max Reach'] = df_prensa_grouped['Max Reach'].apply(format_reach)
    return df_prensa_grouped

def top_influencers_redes_sociales_by_posts(df_cleaned):
    filter_redes_sociales = df_cleaned['Plataforma'] == "Redes Sociales"
    df_redes = df_cleaned[filter_redes_sociales]
    df_redes_grouped = (
        df_redes.groupby('Influencer')[['Reach', 'Source']]
        .agg(Posts=('Reach', 'count'), Max_Reach=('Reach', 'max'), Source=('Source', 'first'))
    )
    df_redes_grouped.reset_index(inplace=True)  # Asegúrate de que 'Influencer' sea una columna normal
    df_redes_grouped['Max_Reach'] = df_redes_grouped['Max_Reach'].apply(format_reach)
    df_redes_grouped = df_redes_grouped.sort_values(by='Posts', ascending=False).head(10)
    return df_redes_grouped[['Influencer', 'Posts', 'Max_Reach', 'Source']]

def top_influencers_redes_sociales_by_reach(df_cleaned):
    filter_redes_sociales = df_cleaned['Plataforma'] == "Redes Sociales"
    df_redes = df_cleaned[filter_redes_sociales]
    df_redes_grouped = df_redes.groupby('Influencer')[['Reach']].agg(['count', 'max'])
    df_redes_grouped.columns = ['Posts', 'Max Reach']
    df_redes_grouped = df_redes_grouped.sort_values(by='Max Reach', ascending=False).head(10)
    df_redes_grouped.reset_index(inplace=True)  # Asegúrate de que 'Influencer' sea una columna normal
    df_redes_grouped['Max Reach'] = df_redes_grouped['Max Reach'].apply(format_reach)
    return df_redes_grouped[['Influencer', 'Max Reach', 'Posts']]

# Función para añadir un DataFrame como tabla en una diapositiva
def add_dataframe_as_table(slide, dataframe, x, y, width, height):
    if dataframe.empty:
        print("DataFrame está vacío.")
        return
    
    # Número de filas y columnas
    rows, cols = dataframe.shape
    table = slide.shapes.add_table(rows + 1, cols, x, y, width, height).table

    # Establecer encabezados de columna
    for j in range(cols):
        cell = table.cell(0, j)
        cell.text = str(dataframe.columns[j])
        
        # Establecer color de fondo del encabezado
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 192, 0)  # Color de fondo del encabezado
        
        # Aplicar formato al texto del encabezado
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True  # Encabezados en negrita
            paragraph.font.size = Pt(14)  # Tamaño de letra
            paragraph.font.name = 'Effra'  # Tipo de letra
            paragraph.alignment = PP_ALIGN.CENTER  # Centrar el texto
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # Centrar verticalmente

    # Llenar la tabla con los datos del DataFrame
    for i in range(rows):
        for j in range(cols):
            value = dataframe.iat[i, j]
            cell = table.cell(i + 1, j)
            cell.text = str(value)
            
            # Centrar texto en las filas
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)  # Tamaño de letra para las filas
                paragraph.font.name = 'Effra Light'  # Tipo de letra para las filas
                paragraph.alignment = PP_ALIGN.CENTER  # Centrar texto horizontalmente
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # Centrar verticalmente

    print("Tabla añadida correctamente.")

            
# Función para guardar el dataframe modificado como archivo csv
def save_cleaned_csv(df, file_path):
    base_filename = os.path.basename(file_path).split()[0]
    base_filename = base_filename.split('.')[0]
    output_filename = f"{base_filename}_(resultado).csv"
    with open(output_filename, 'w', encoding='utf-16', newline='') as f:
        df.to_csv(f, index=False, sep='\t')
    print(f"CSV data saved to: {output_filename}")