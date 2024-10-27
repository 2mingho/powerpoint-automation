import os
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
import calculation as report

# Variables y rutas de archivo
file_path = "scratch/MTRABAJO - Oct 14, 2024 - 12 01 11 PM.csv"
pptx_template_path = "powerpoints/Reporte_plantilla.pptx"

# Cargar y limpiar datos
df_cleaned = report.load_and_clean_data(file_path)
df_cleaned['Influencer'] = df_cleaned.apply(report.update_influencer, axis=1)
df_cleaned['Sentiment'] = df_cleaned.apply(report.update_sentiment, axis=1)

# Obtener métricas
total_mentions, count_of_authors, estimated_reach = report.calculate_summary_metrics(df_cleaned)
formatted_estimated_reach = estimated_reach

# Guardar datos en csv
report.save_cleaned_csv(df_cleaned, file_path)

# Crear gráficos
report.create_mentions_evolution_chart(df_cleaned)
report.create_sentiment_pie_chart(df_cleaned)

# Obtener distribución de plataformas
platform_counts, max_reach_per_platform = report.distribucion_plataforma(df_cleaned)

# Obtener las noticias más relevantes
top_sentences = report.get_top_hit_sentences(df_cleaned)

# Obtener influencers
top_influencers_prensa = report.top_influencers_prensa_digital(df_cleaned)
top_influencers_redes_posts = report.top_influencers_redes_sociales_by_posts(df_cleaned)
top_influencers_redes_reach = report.top_influencers_redes_sociales_by_reach(df_cleaned)

# Obtener nombre del archivo y cliente
current_date = datetime.now().strftime('%d-%b-%Y')
current_date_file_name = datetime.now().strftime('%d-%b-%Y, %H %M %S')
client_name = os.path.basename(file_path).split()[0]

# Cargar presentación y modificar slides
prs = Presentation(pptx_template_path)

# Modificar slide 1
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame:
        for key, value in {"REPORT_CLIENT": client_name, "REPORT_DATE": current_date}.items():
            if key in shape.text:
                report.set_text_style(shape, str(value), 'Effra Heavy', Pt(28), False)

# Modificar slide 2
slide2 = prs.slides[1]
for shape in slide2.shapes:
    if shape.has_text_frame:
        for key, value in {
            "NUMB_MENTIONS": str(total_mentions),
            "NUMB_ACTORS": str(count_of_authors),
            "EST_REACH": formatted_estimated_reach
        }.items():
            if key in shape.text:
                report.set_text_style(shape, str(value), font_size=Pt(22), center=True)

    elif shape.shape_type == 13:
        slide2.shapes.add_picture('scratch/convEvolution.png', Inches(0.2), Inches(1.2), width=Inches(10.46), height=Inches(5.63))

# Modificar slide 3
slide3 = prs.slides[2]
for shape in slide3.shapes:
    if shape.has_text_frame:
        for key, value in {"TOP_NEWS": "\n".join(top_sentences)}.items():
            if key in shape.text:
                report.set_text_style(shape, str(value), 'Effra Light', Pt(12), False)
try:
    slide3.shapes.add_picture('scratch/Wordcloud.png', Inches(7.5), Inches(3.5), width=Inches(4.2), height=Inches(2.66)) 
except Exception as e:
    print(e)

# Modificar slide 4
slide4 = prs.slides[3]
for shape in slide4.shapes:
    if shape.shape_type == 13:
        slide4.shapes.add_picture('scratch/sentiment_pie_chart.png', Inches(1), Inches(1), width=Inches(6), height=Inches(6))

# Modificar slide 5
slide5 = prs.slides[4]
for shape in slide5.shapes:
    if shape.has_text_frame:
        for key, value in {"NUMB_PRENSA": platform_counts['Prensa Digital']}.items():
            if key in shape.text:
                report.set_text_style(shape, str(value), font_size=Pt(28))
try:
    report.add_dataframe_as_table(slide5, top_influencers_prensa, Inches(2.65), Inches(2), Inches(8), Inches(4))
except Exception as e:
    print("Error al añadir la tabla en slide5:", e)

# Modificar slide 6
slide6 = prs.slides[5]
for shape in slide6.shapes:
    if shape.has_text_frame:
        for key, value in {"NUMB_REDES": platform_counts['Redes Sociales']}.items():
            if key in shape.text:
                report.set_text_style(shape, str(value), font_size=Pt(28))
try:
    report.add_dataframe_as_table(slide6, top_influencers_redes_posts, Inches(0.56), Inches(2), Inches(7), Inches(4))
except Exception as e:
    print("Error al añadir la tabla en slide6 (posts):", e)

try:
    report.add_dataframe_as_table(slide6, top_influencers_redes_reach, Inches(8), Inches(2), Inches(5), Inches(4))
except Exception as e:
    print("Error al añadir la tabla en slide6 (reach):", e)

# Guardar presentación
output_filename = f"{client_name} {current_date_file_name}.pptx"
prs.save(output_filename)
print(f"Presentación guardada como {output_filename}")
